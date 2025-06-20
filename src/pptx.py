from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
import os

def create_presentation(output_file="demo.pptx"):
    """创建包含多种元素的PPT演示文稿"""
    
    # 初始化演示文稿
    prs = Presentation()
    
    # ==================== 1. 添加标题幻灯片 ====================
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片版式
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "自动化PPT演示"
    subtitle.text = "使用python-pptx创建\n包含文本、图片、表格和图表"
    
    # ==================== 2. 添加文本幻灯片 ====================
    slide_layout = prs.slide_layouts[1]  # 标题和内容版式
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "文本内容示例"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # 添加段落
    p = tf.add_paragraph()
    p.text = "这是第一级文本"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "这是第二级文本"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "这是第三级文本，可以设置不同的颜色和样式"
    p.level = 2
    p.font.color.rgb = RGBColor(0x42, 0x24, 0xE9)  # 蓝色
    
    # ==================== 3. 添加图片幻灯片 ====================
    slide_layout = prs.slide_layouts[5]  # 仅标题版式
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "图片示例"
    
    # 添加图片 (确保当前目录有sample.jpg或替换为你的图片路径)
    img_path = "sample.jpg"
    if not os.path.exists(img_path):
        print(f"警告: {img_path} 不存在，将创建空白图片占位符")
        left = Inches(1)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        slide.shapes.add_shape(
            1, left, top, width, height
        ).text = "图片占位符"
    else:
        left = Inches(1)
        top = Inches(2)
        height = Inches(4)
        pic = slide.shapes.add_picture(
            img_path, left, top, height=height
        )
        
        # 添加图片说明
        left = Inches(1)
        top = Inches(6.2)
        width = Inches(6)
        height = Inches(0.5)
        txBox = slide.sh
