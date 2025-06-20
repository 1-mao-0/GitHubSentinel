from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from enum import Enum, auto
import sys
import os
from typing import Dict, List, Optional

class ContentType(Enum):
    """幻灯片内容类型枚举"""
    TITLE = auto()
    SECTION = auto()
    TEXT = auto()
    IMAGE = auto()
    CHART = auto()
    TABLE = auto()
    MIXED = auto()
    QUOTE = auto()
    END = auto()

class PPTLayoutConverter:
    """PPT布局转换核心类"""
    
    def __init__(self):
        self.layout_library: Dict[ContentType, List] = {}
        self.stats = {
            'total_slides': 0,
            'converted': 0,
            'warnings': 0
        }
    
    def analyze_content(self, slide) -> ContentType:
        """分析幻灯片内容类型"""
        shape_stats = {
            'text': 0,
            'image': 0,
            'chart': 0,
            'table': 0
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_stats['text'] += 1
            if shape.shape_type == 13:  # 图片类型
                shape_stats['image'] += 1
            if shape.has_chart:
                shape_stats['chart'] += 1
            if shape.has_table:
                shape_stats['table'] += 1
        
        # 首页/末页判断
        if self.stats['total_slides'] == 0:
            return ContentType.TITLE
        if slide == slide.presentation.slides[-1]:
            return ContentType.END
        
        # 内容类型判断逻辑
        if shape_stats['chart'] > 0:
            return ContentType.CHART
        if shape_stats['table'] > 0:
            return ContentType.TABLE
        if shape_stats['image'] > 1 and shape_stats['text'] > 0:
            return ContentType.MIXED
        if shape_stats['image'] == 1 and shape_stats['text'] < 3:
            return ContentType.IMAGE
        if shape_stats['text'] > 5:
            return ContentType.TEXT
        if "引用" in slide.shapes[0].text or "quote" in slide.shapes[0].text.lower():
            return ContentType.QUOTE
        
        return ContentType.SECTION
    
    def collect_layouts(self, presentation) -> None:
        """从模板收集各类版式"""
        for slide in presentation.slides:
            content_type = self.analyze_content(slide)
            if content_type not in self.layout_library:
                self.layout_library[content_type] = []
            self.layout_library[content_type].append(slide.slide_layout)
    
    def add_verification_mark(self, slide) -> None:
        """添加版式验证标记"""
        left = Inches(8.5)
        top = Inches(6.7)
        width = Inches(1.5)
        height = Inches(0.4)
        
        # 避免覆盖现有内容
        for shape in slide.shapes:
            if shape.left > left and shape.top > top:
                left = shape.left - width - Inches(0.2)
                top = shape.top - height - Inches(0.2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"版式: {slide.slide_layout.name}"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(255, 0, 0)
    
    def check_content_fit(self, slide) -> bool:
        """检查内容是否适应新版式"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_length = len(shape.text)
                if text_length > 1000:  # 长文本警告
                    print(f"警告：幻灯片{slide.slide_id}可能文本过长 ({text_length}字符)")
                    self.stats['warnings'] += 1
                    return False
                if shape.text_frame.overflow:
                    print(f"警告：幻灯片{slide.slide_id}文本溢出")
                    self.stats['warnings'] += 1
                    return False
        return True
    
    def convert_presentation(self, source_file: str, target_file: str, output_file: str) -> bool:
        """执行PPT转换主函数"""
        try:
            print("正在加载源模板文件...")
            source_ppt = Presentation(source_file)
            
            print("正在分析模板版式...")
            self.collect_layouts(source_ppt)
            
            print("正在加载目标PPT...")
            target_ppt = Presentation(target_file)
            self.stats['total_slides'] = len(target_ppt.slides)
            
            print("开始版式转换...")
            for i, slide in enumerate(target_ppt.slides):
                content_type = self.analyze_content(slide)
                
                if content_type in self.layout_library:
                    # 选择最匹配的版式
                    selected_layout = self.layout_library[content_type][i % len(self.layout_library[content_type])]
                    slide.slide_layout = selected_layout
                    self.stats['converted'] += 1
                    
                    # 添加验证标记
                    self.add_verification_mark(slide)
                    
                    # 检查内容适配
                    if not self.check_content_fit(slide):
                        print(f"请注意检查幻灯片 {i+1} 的内容显示")
                
                print(f"进度: {i+1}/{self.stats['total_slides']}", end='\r')
            
            print("\n正在保存转换结果...")
            target_ppt.save(output_file)
            
            print("转换完成！")
            print(f"统计信息:")
            print(f"- 总幻灯片数: {self.stats['total_slides']}")
            print(f"- 成功转换: {self.stats['converted']}")
            print(f"- 需要检查的幻灯片: {self.stats['warnings']}")
            
            return True
            
        except Exception as e:
            print(f"转换过程中出错: {str(e)}")
            return False

def main():
    if len(sys.argv) != 4:
        print("PPT布局转换工具 v1.0")
        print("用法: python ppt_converter.py <模板文件> <输入文件> <输出文件>")
        print("示例: python ppt_converter.py FairFrames.pptx ChatPPT.pptx Output.pptx")
        return
    
    if not os.path.exists(sys.argv[1]):
        print(f"错误: 模板文件 {sys.argv[1]} 不存在")
        return
    
    if not os.path.exists(sys.argv[2]):
        print(f"错误: 输入文件 {sys.argv[2]} 不存在")
        return
    
    converter = PPTLayoutConverter()
    if converter.convert_presentation(sys.argv[1], sys.argv[2], sys.argv[3]):
        print(f"结果已保存到 {sys.argv[3]}")
        print("请打开PPT检查红色版式标记，确认内容显示正常")

if __name__ == "__main__":
    main()
